import os
import google.generativeai as genai
import chromadb
from pypdf import PdfReader
import io
import docx
import textract
import uuid
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# Initialize the embedding model
_embedding_model = None

def get_embedding_model():
    """Get or initialize the embedding model"""
    global _embedding_model
    if _embedding_model is None:
        # Using a lightweight model - you can replace with others like 'all-mpnet-base-v2' for better quality
        _embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    return _embedding_model

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
    return text

def extract_text_from_docx(file_path):
    """Extract text from DOCX files using python-docx"""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception:
        # Fallback to textract if python-docx fails
        return extract_text_with_textract(file_path)

def extract_text_from_txt(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encodings if utf-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception:
            return extract_text_with_textract(file_path)

def extract_text_from_pptx(file_path):
    """Extract text from PPTX files using python-pptx"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        # Fallback to textract if python-pptx fails or for .ppt files
        return extract_text_with_textract(file_path)

def extract_text_with_textract(file_path):
    """Extract text using textract as a fallback method"""
    try:
        return textract.process(file_path).decode('utf-8')
    except Exception as e:
        raise Exception(f"Failed to extract text with textract: {str(e)}")

def extract_text_from_document(file_path):
    """Extract text from a document based on its file extension"""
    _, file_extension = os.path.splitext(file_path.lower())
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.doc' or file_extension == '.rtf':
        return extract_text_with_textract(file_path)
    elif file_extension == '.txt':
        return extract_text_from_txt(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_extension == '.ppt':  # Older .ppt format, try textract
        return extract_text_with_textract(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")

def chunk_text(text, chunk_size=1000, overlap=100):
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size - overlap)]
    return chunks

def get_embeddings(text_chunks):
    """Generate embeddings using local model instead of Google's API"""
    model = get_embedding_model()
    embeddings = model.encode(text_chunks)
    # Return as list of lists to match the expected format
    return embeddings.tolist()

def initialize_chromadb(db_path, collection_name):
    """Initialize ChromaDB with the specified path and collection name
    
    Args:
        db_path: Path to the ChromaDB persistence directory
        collection_name: Name of the collection to use
        
    Returns:
        A ChromaDB collection object
    """
    try:
        client = chromadb.PersistentClient(path=db_path)
        return client.get_or_create_collection(name=collection_name)
    except Exception as e:
        print(f"Error initializing ChromaDB: {e}")
        return None

def query_chromadb(collection, query_text, n_results=5, document_names=None):
    """Query ChromaDB using local embedding model"""
    try:
        model = get_embedding_model()
        query_embedding = model.encode(query_text).tolist()
        
        query_params = {
            "query_embeddings": [query_embedding],
            "n_results": n_results,
            "include": ['documents', 'metadatas']
        }

        if document_names and isinstance(document_names, list) and len(document_names) > 0:
            query_params["where"] = {"filename": {"$in": document_names}}

        results = collection.query(**query_params)
        return results.get('documents', [[]])[0]
    except Exception as e:
        print(f"Error querying ChromaDB: {e}")
        return []

def generate_response(query_text, context_chunks):
    """Generate a response using Google Gemini model
    
    Make sure GEMINI_API_KEY is set in your environment or Flask config
    """
    try:
        context = "\n\n".join(context_chunks)
        prompt = f"""Based ONLY on the following context:\n\n{context}\n\nAnswer the following question:\n{query_text}\n\nAnswer:"""
        response = genai.GenerativeModel('gemini-1.5-flash').generate_content(prompt)
        return response.text if response.text else "No answer available."
    except Exception as e:
        print(f"Error generating response: {e}")
        return "An error occurred while generating the response."
